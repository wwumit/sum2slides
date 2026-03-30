#!/usr/bin/env python3
"""
生成Sum2Slides v1.0.1典型场景演示幻灯片
专门展示"讨论总结转PPT"这个最常用场景
"""

import os
from datetime import datetime
from pptx import Presentation

def create_typical_scenario():
    """创建典型场景演示幻灯片"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 幻灯片1：标题页 - 突出典型场景
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Sum2Slides v1.0.1"
    subtitle.text = "典型使用场景\n讨论总结转PPT\n" + datetime.now().strftime("%Y-%m-%d %H:%M")
    
    # 幻灯片2：典型场景描述
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "🎯 最常用场景"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "场景描述："
    p = text_frame.add_paragraph()
    p.text = '"请将上述讨论中，对主要事件的讨论过程和结论，形成一个总结性的PPT"'
    
    p = text_frame.add_paragraph()
    p.text = "🔍 这个场景的特点："
    p = text_frame.add_paragraph()
    p.text = "• 团队讨论后的快速总结"
    p = text_frame.add_paragraph()
    p.text = "• 自动结构化呈现讨论内容"
    p = text_frame.add_paragraph()
    p.text = "• 便于汇报和后续跟进"
    p = text_frame.add_paragraph()
    p.text = "• 节省大量手动整理时间"
    
    # 幻灯片3：实际使用示例 - 飞书群聊场景
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    
    title.text = "💬 实际使用示例"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "在飞书群聊中的使用："
    p = text_frame.add_paragraph()
    p.text = "```"
    p = text_frame.add_paragraph()
    p.text = "【团队讨论中...】"
    p = text_frame.add_paragraph()
    p.text = "张三：我们讨论了项目延期的问题"
    p = text_frame.add_paragraph()
    p.text = "李四：主要原因有需求变更和资源不足"
    p = text_frame.add_paragraph()
    p.text = "王五：解决方案是增加人手和调整排期"
    p = text_frame.add_paragraph()
    p.text = "赵六：需要向领导汇报这个情况"
    p = text_frame.add_paragraph()
    p.text = ""
    p = text_frame.add_paragraph()
    p.text = "经理：@sum2slides 请将上述讨论中，对主要事件的讨论过程和结论，形成一个总结性的PPT"
    p = text_frame.add_paragraph()
    p.text = "```"
    
    # 幻灯片4：sum2slides的处理过程
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    content = slide4.placeholders[1]
    
    title.text = "🔄 sum2slides 处理过程"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "1. 📥 接收指令"
    p = text_frame.add_paragraph()
    p.text = "   识别关键词：'讨论中'、'主要事件'、'讨论过程'、'结论'"
    
    p = text_frame.add_paragraph()
    p.text = "2. 🔍 分析讨论内容"
    p = text_frame.add_paragraph()
    p.text = "   • 提取主要事件：项目延期"
    p = text_frame.add_paragraph()
    p.text = "   • 识别讨论过程：原因分析"
    p = text_frame.add_paragraph()
    p.text = "   • 提取结论：解决方案"
    
    p = text_frame.add_paragraph()
    p.text = "3. 🎨 智能生成结构"
    p = text_frame.add_paragraph()
    p.text = "   自动创建："
    p = text_frame.add_paragraph()
    p.text = "   • 封面页"
    p = text_frame.add_paragraph()
    p.text = "   • 问题概述"
    p = text_frame.add_paragraph()
    p.text = "   • 讨论过程"
    p = text_frame.add_paragraph()
    p.text = "   • 原因分析"
    p = text_frame.add_paragraph()
    p.text = "   • 解决方案"
    p = text_frame.add_paragraph()
    p.text = "   • 行动计划"
    
    # 幻灯片5：生成的PPT内容预览
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    content = slide5.placeholders[1]
    
    title.text = "📄 生成的PPT内容"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "生成的PPT包含以下部分："
    p = text_frame.add_paragraph()
    p.text = ""
    p = text_frame.add_paragraph()
    p.text = "📋 第1页：封面"
    p = text_frame.add_paragraph()
    p.text = "   • 标题：项目延期问题讨论总结"
    p = text_frame.add_paragraph()
    p.text = "   • 副标题：团队讨论纪要"
    p = text_frame.add_paragraph()
    p.text = "   • 日期：2026-03-29"
    
    p = text_frame.add_paragraph()
    p.text = "📋 第2页：问题概述"
    p = text_frame.add_paragraph()
    p.text = "   • 主要事件：项目延期两周"
    p = text_frame.add_paragraph()
    p.text = "   • 影响范围：版本发布延迟"
    p = text_frame.add_paragraph()
    p.text = "   • 紧急程度：高"
    
    p = text_frame.add_paragraph()
    p.text = "📋 第3页：讨论过程"
    p = text_frame.add_paragraph()
    p.text = "   • 团队成员观点汇总"
    p = text_frame.add_paragraph()
    p.text = "   • 问题分析过程"
    p = text_frame.add_paragraph()
    p.text = "   • 关键讨论要点"
    
    # 幻灯片6：使用方式总结
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    content = slide6.placeholders[1]
    
    title.text = "🚀 使用方式总结"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "💡 触发方式："
    p = text_frame.add_paragraph()
    p.text = "1. @sum2slides + 指令"
    p = text_frame.add_paragraph()
    p.text = "   ```"
    p = text_frame.add_paragraph()
    p.text = "   @sum2slides 请将上述讨论总结为PPT"
    p = text_frame.add_paragraph()
    p.text = "   ```"
    
    p = text_frame.add_paragraph()
    p.text = "2. /sum2slides + 参数"
    p = text_frame.add_paragraph()
    p.text = "   ```"
    p = text_frame.add_paragraph()
    p.text = "   /sum2slides --type summary --template business"
    p = text_frame.add_paragraph()
    p.text = "   ```"
    
    p = text_frame.add_paragraph()
    p.text = "3. 关键词自动触发"
    p = text_frame.add_paragraph()
    p.text = "   • '总结性PPT'"
    p = text_frame.add_paragraph()
    p.text = "   • '讨论纪要'"
    p = text_frame.add_paragraph()
    p.text = "   • '汇报材料'"
    
    # 幻灯片7：实际应用价值
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    content = slide7.placeholders[1]
    
    title.text = "💰 实际应用价值"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "✅ 时间节省"
    p = text_frame.add_paragraph()
    p.text = "   • 手动制作：2-3小时"
    p = text_frame.add_paragraph()
    p.text = "   • sum2slides：2-3分钟"
    p = text_frame.add_paragraph()
    p.text = "   • 效率提升：60倍+"
    
    p = text_frame.add_paragraph()
    p.text = "✅ 质量保证"
    p = text_frame.add_paragraph()
    p.text = "   • 结构化呈现"
    p = text_frame.add_paragraph()
    p.text = "   • 专业模板"
    p = text_frame.add_paragraph()
    p.text = "   • 一致性格式"
    
    p = text_frame.add_paragraph()
    p.text = "✅ 团队协作"
    p = text_frame.add_paragraph()
    p.text = "   • 快速分享"
    p = text_frame.add_paragraph()
    p.text = "   • 统一认知"
    p = text_frame.add_paragraph()
    p.text = "   • 便于跟进"
    
    # 幻灯片8：更多指令示例
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    content = slide8.placeholders[1]
    
    title.text = "💬 更多指令示例"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "场景1：会议总结"
    p = text_frame.add_paragraph()
    p.text = "```"
    p = text_frame.add_paragraph()
    p.text = "@sum2slides 将刚才的会议讨论要点整理成PPT，用于部门汇报"
    p = text_frame.add_paragraph()
    p.text = "```"
    
    p = text_frame.add_paragraph()
    p.text = "场景2：项目复盘"
    p = text_frame.add_paragraph()
    p.text = "```"
    p = text_frame.add_paragraph()
    p.text = "@sum2slides 请把项目复盘的主要发现和建议做成演示文稿"
    p = text_frame.add_paragraph()
    p.text = "```"
    
    p = text_frame.add_paragraph()
    p.text = "场景3：学习分享"
    p = text_frame.add_paragraph()
    p.text = "```"
    p = text_frame.add_paragraph()
    p.text = "@sum2slides 把这次技术分享的内容转为PPT，方便后续培训使用"
    p = text_frame.add_paragraph()
    p.text = "```"
    
    # 幻灯片9：最佳实践建议
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide9.shapes.title
    content = slide9.placeholders[1]
    
    title.text = "🌟 最佳实践建议"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "1. 🎯 明确指令"
    p = text_frame.add_paragraph()
    p.text = "   • 说明用途：'用于汇报'、'用于分享'"
    p = text_frame.add_paragraph()
    p.text = "   • 指定受众：'给领导看'、'给团队用'"
    p = text_frame.add_paragraph()
    p.text = "   • 明确重点：'突出解决方案'、'强调数据'"
    
    p = text_frame.add_paragraph()
    p.text = "2. 📝 提供上下文"
    p = text_frame.add_paragraph()
    p.text = "   • 包含讨论背景"
    p = text_frame.add_paragraph()
    p.text = "   • 说明关键决策"
    p = text_frame.add_paragraph()
    p.text = "   • 提供相关数据"
    
    p = text_frame.add_paragraph()
    p.text = "3. 🎨 利用模板优势"
    p = text_frame.add_paragraph()
    p.text = "   • 商务场景：business模板"
    p = text_frame.add_paragraph()
    p.text = "   • 学术场景：academic模板"
    p = text_frame.add_paragraph()
    p.text = "   • 创意场景：colorful模板"
    
    # 幻灯片10：总结
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide10.shapes.title
    subtitle = slide10.placeholders[1]
    
    title.text = "Sum2Slides 让讨论总结更高效"
    subtitle.text = "💬 讨论过程 → 📊 结构化PPT\n\n" + \
                   "🚀 一键转换，专业呈现\n" + \
                   "🤝 团队协作，知识沉淀\n" + \
                   "💡 智能分析，精准总结\n\n" + \
                   "版本: Sum2Slides v1.0.1\n" + \
                   "典型场景: 讨论总结转PPT"
    
    # 保存文件
    output_dir = "typical_scenario"
    os.makedirs(output_dir, exist_ok=True)
    
    output_file = f"{output_dir}/sum2slides_typical_scenario.pptx"
    prs.save(output_file)
    
    print(f"✅ 典型场景演示幻灯片已生成: {output_file}")
    print(f"📊 幻灯片数量: {len(prs.slides)} 张")
    
    return output_file

def main():
    """主函数"""
    
    print("🎬 生成Sum2Slides v1.0.1典型场景演示幻灯片")
    print("=" * 50)
    print("典型场景：'请将上述讨论中，对主要事件的讨论过程和结论，形成一个总结性的PPT'")
    print("=" * 50)
    
    try:
        # 检查python-pptx是否安装
        import pptx
        print("✅ python-pptx 已安装")
        
        # 生成演示文稿
        output_file = create_typical_scenario()
        
        print("\n✅ 典型场景演示幻灯片生成完成！")
        print(f"\n📁 输出文件: {output_file}")
        print("\n📋 幻灯片内容:")
        print("  1. 标题页 - 典型使用场景演示")
        print("  2. 最常用场景介绍")
        print("  3. 实际对话示例")
        print("  4. sum2slides处理过程")
        print("  5. 生成的PPT内容")
        print("  6. 使用方式总结")
        print("  7. 实际应用价值")
        print("  8. 更多指令示例")
        print("  9. 最佳实践建议")
        print("  10. 总结")
        
        print("\n🎯 重点展示:")
        print("  • 实际对话示例：飞书群聊中的真实使用")
        print("  • 处理过程：如何智能分析讨论内容")
        print("  • 生成结果：结构化的PPT内容预览")
        print("  • 多种触发方式：@sum2slides、/sum2slides、关键词")
        
        # 复制到samba共享文件夹
        samba_path = "~/samba_shares/小七文档/skill-versions/sum2slides_typical_scenario.pptx"
        import shutil
        shutil.copy(output_file, os.path.expanduser(samba_path))
        
        print(f"\n📤 已复制到Samba共享文件夹: {samba_path}")
        
    except ImportError as e:
        print(f"❌ 需要安装python-pptx: {e}")
        print("\n💡 安装方法:")
        print("  pip install python-pptx")

if __name__ == "__main__":
    main()