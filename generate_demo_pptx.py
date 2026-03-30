#!/usr/bin/env python3
"""
生成Sum2Slides v1.0.1演示幻灯片
使用python-pptx直接创建演示文稿，展示功能和用法
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_demo_presentation():
    """创建演示幻灯片"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 幻灯片1：标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Sum2Slides v1.0.1"
    subtitle.text = "正式版本演示\n" + datetime.now().strftime("%Y-%m-%d %H:%M")
    
    # 幻灯片2：简介
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "什么是Sum2Slides？"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🚀 文本转幻灯片工具"
    p = text_frame.add_paragraph()
    p.text = "🎯 自动生成专业演示文稿"
    p = text_frame.add_paragraph()
    p.text = "📊 支持PPTX格式输出"
    p = text_frame.add_paragraph()
    p.text = "💡 智能内容分析"
    
    # 幻灯片3：核心价值
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    
    title.text = "核心价值"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "✅ 节省时间 - 自动转换，减少手动工作"
    p = text_frame.add_paragraph()
    p.text = "✅ 保持一致性 - 统一格式和风格"
    p = text_frame.add_paragraph()
    p.text = "✅ 提高质量 - 结构化展示，逻辑清晰"
    p = text_frame.add_paragraph()
    p.text = "✅ 易于修改 - 基于模板，快速调整"
    
    # 幻灯片4：安装方式
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    content = slide4.placeholders[1]
    
    title.text = "安装方式"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "📦 通过ClawHub安装（推荐）"
    p = text_frame.add_paragraph()
    p.text = "```bash\nopenclaw skill install sum2slides\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🐍 通过pip安装"
    p = text_frame.add_paragraph()
    p.text = "```bash\npip install sum2slides\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🔧 从源码安装"
    p = text_frame.add_paragraph()
    p.text = "```bash\n通过 OpenClaw 安装技能\ncd sum2slides\npip install -e .\n```"
    
    # 幻灯片5：基本使用
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    content = slide5.placeholders[1]
    
    title.text = "基本使用"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "📝 命令行使用"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 从文本生成幻灯片\nsum2slides \"你的文本内容\" --output demo.pptx\n\n# 从文件输入\nsum2slides --input notes.txt --output slides.pptx\n\n# 使用特定模板\nsum2slides --input summary.md --template business --output report.pptx\n```"
    
    # 幻灯片6：Python API
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    content = slide6.placeholders[1]
    
    title.text = "Python API 使用"
    text_frame = content.text_frame
    text_frame.clear()
    
    code = """from sum2slides import Sum2Slides

# 创建转换器
converter = Sum2Slides()

# 转换文本
text = \"\"\"项目更新

本周完成工作：
1. 完成用户认证模块
2. 优化数据库性能
3. 修复已知问题

下周计划：
1. 开发新功能
2. 进行系统测试
\"\"\"

presentation = converter.convert(text)
presentation.save("output.pptx")"""
    
    p = text_frame.paragraphs[0]
    p.text = code
    
    # 幻灯片7：高级功能
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    content = slide7.placeholders[1]
    
    title.text = "高级功能"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🎨 模板系统"
    p = text_frame.add_paragraph()
    p.text = "  • 默认模板 - 通用场景"
    p = text_frame.add_paragraph()
    p.text = "  • 商务模板 - 商务汇报"
    p = text_frame.add_paragraph()
    p.text = "  • 学术模板 - 学术报告"
    
    p = text_frame.add_paragraph()
    p.text = "🌞 主题支持"
    p = text_frame.add_paragraph()
    p.text = "  • 浅色主题 - 适合明亮环境"
    p = text_frame.add_paragraph()
    p.text = "  • 深色主题 - 适合暗光环境"
    p = text_frame.add_paragraph()
    p.text = "  • 彩色主题 - 创意展示"
    
    # 幻灯片8：使用场景
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    content = slide8.placeholders[1]
    
    title.text = "使用场景"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "📊 会议纪要转演示"
    p = text_frame.add_paragraph()
    p.text = "```bash\nsum2slides --input meeting_notes.txt --template business --output meeting_summary.pptx\n```"
    
    p = text_frame.add_paragraph()
    p.text = "💼 项目报告生成"
    p = text_frame.add_paragraph()
    p.text = "```bash\nsum2slides --input project_report.md --template business --output project_presentation.pptx\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🎓 学术汇报制作"
    p = text_frame.add_paragraph()
    p.text = "```bash\nsum2slides --input research_paper.md --template academic --output research_presentation.pptx\n```"
    
    # 幻灯片9：最佳实践
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide9.shapes.title
    content = slide9.placeholders[1]
    
    title.text = "最佳实践"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "📝 文本格式建议"
    p = text_frame.add_paragraph()
    p.text = "```\n# 主标题\n\n## 副标题或章节\n\n• 项目1：详细描述\n• 项目2：详细描述\n\n1. 编号项目1：详细说明\n2. 编号项目2：详细说明\n\n---\n\n## 新章节\n更多内容...\n```"
    
    # 幻灯片10：资源和支持
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide10.shapes.title
    content = slide10.placeholders[1]
    
    title.text = "资源和支持"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "📚 文档资源"
    p = text_frame.add_paragraph()
    p.text = "  • 用户手册: README.md"
    p = text_frame.add_paragraph()
    p.text = "  • 新手教程: TUTORIAL.md"
    p = text_frame.add_paragraph()
    p.text = "  • 常见问题: FAQ.md"
    p = text_frame.add_paragraph()
    p.text = "  • 使用演示: DEMO.md"
    
    p = text_frame.add_paragraph()
    p.text = "🔗 社区支持"
    p = text_frame.add_paragraph()
    p.text = "  • GitHub: 本地化技能版本"
    p = text_frame.add_paragraph()
    p.text = "  • Issues: 本地化技能版本/issues"
    p = text_frame.add_paragraph()
    p.text = "  • Discord: https://discord.gg/clawd"
    
    # 幻灯片11：总结
    slide11 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide11.shapes.title
    subtitle = slide11.placeholders[1]
    
    title.text = "感谢使用Sum2Slides！"
    subtitle.text = "🚀 Happy Presenting!\n\n版本: Sum2Slides v1.0.1\n日期: " + datetime.now().strftime("%Y-%m-%d")
    
    # 保存文件
    output_dir = "generated_demo"
    os.makedirs(output_dir, exist_ok=True)
    
    output_file = f"{output_dir}/sum2slides_demo_presentation.pptx"
    prs.save(output_file)
    
    print(f"✅ 演示幻灯片已生成: {output_file}")
    print(f"📊 幻灯片数量: {len(prs.slides)} 张")
    
    return output_file

def main():
    """主函数"""
    
    print("🎬 生成Sum2Slides v1.0.1演示幻灯片")
    print("=" * 50)
    
    try:
        # 检查python-pptx是否安装
        import pptx
        print("✅ python-pptx 已安装")
        
        # 生成演示文稿
        output_file = create_demo_presentation()
        
        print("\n✅ 演示幻灯片生成完成！")
        print(f"\n📁 输出文件: {output_file}")
        print("\n📋 幻灯片内容:")
        print("  1. 标题页 - Sum2Slides v1.0.1")
        print("  2. 简介 - 什么是Sum2Slides？")
        print("  3. 核心价值 - 节省时间、保持一致性等")
        print("  4. 安装方式 - 3种安装方法")
        print("  5. 基本使用 - 命令行使用")
        print("  6. Python API - 代码示例")
        print("  7. 高级功能 - 模板和主题")
        print("  8. 使用场景 - 实际应用案例")
        print("  9. 最佳实践 - 文本格式建议")
        print("  10. 资源和支持 - 文档和社区")
        print("  11. 总结 - 感谢使用")
        
        print("\n🎯 您可以：")
        print("  1. 使用PowerPoint或LibreOffice打开查看")
        print("  2. 检查使用说明是否清晰")
        print("  3. 验证幻灯片效果")
        
    except ImportError as e:
        print(f"❌ 需要安装python-pptx: {e}")
        print("\n💡 安装方法:")
        print("  pip install python-pptx")
        print("\n或者，让我创建一个简单的演示内容文件...")
        
        # 创建简单的演示内容
        simple_demo = """# Sum2Slides v1.0.1 演示内容

## 由于python-pptx未安装，无法生成PPTX文件

### 请安装依赖：
```bash
pip install python-pptx
```

### 然后重新运行生成脚本

### 演示内容预览：

1. **标题页**: Sum2Slides v1.0.1
2. **简介**: 文本转幻灯片工具
3. **核心价值**: 节省时间、保持一致性等
4. **安装方式**: 3种安装方法
5. **基本使用**: 命令行示例
6. **Python API**: 代码示例
7. **高级功能**: 模板和主题
8. **使用场景**: 实际应用
9. **最佳实践**: 格式建议
10. **资源支持**: 文档和社区
11. **总结**: 感谢使用
"""
        
        output_dir = "generated_demo"
        os.makedirs(output_dir, exist_ok=True)
        
        text_file = f"{output_dir}/sum2slides_demo_content.txt"
        with open(text_file, "w", encoding="utf-8") as f:
            f.write(simple_demo)
        
        print(f"\n📝 已创建演示内容文件: {text_file}")
        print("💡 您可以查看此文件了解演示内容")

if __name__ == "__main__":
    main()