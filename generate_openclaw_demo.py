#!/usr/bin/env python3
"""
生成Sum2Slides v1.0.1演示幻灯片（包含OpenClaw使用说明）
专门展示在OpenClaw场景下的使用方法
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_openclaw_demo_presentation():
    """创建包含OpenClaw使用说明的演示幻灯片"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 幻灯片1：标题页（强调OpenClaw Skill）
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Sum2Slides v1.0.1"
    subtitle.text = "OpenClaw Skill 演示\n专为OpenClaw场景设计\n" + datetime.now().strftime("%Y-%m-%d %H:%M")
    
    # 幻灯片2：OpenClaw Skill简介
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "OpenClaw Skill 定位"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🎯 专为OpenClaw设计的文本转幻灯片工具"
    p = text_frame.add_paragraph()
    p.text = "🔧 深度集成OpenClaw技能生态系统"
    p = text_frame.add_paragraph()
    p.text = "🚀 提供一键式安装和配置"
    p = text_frame.add_paragraph()
    p.text = "💡 支持OpenClaw的所有使用场景"
    
    # 幻灯片3：OpenClaw核心安装方式
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    
    title.text = "OpenClaw Skill 安装"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🎯 推荐方式：通过ClawHub安装"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 一键安装Sum2Slides Skill\nopenclaw skill install sum2slides\n```"
    
    p = text_frame.add_paragraph()
    p.text = "📦 安装后验证"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 验证Skill安装成功\nopenclaw skill list | grep sum2slides\n\n# 查看Skill详情\nopenclaw skill info sum2slides\n```"
    
    # 幻灯片4：OpenClaw环境下的使用
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    content = slide4.placeholders[1]
    
    title.text = "OpenClaw环境使用"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🔧 在OpenClaw会话中直接使用"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 在OpenClaw聊天中调用\n@sum2slides 将会议纪要转为幻灯片\n\n# 或使用命令模式\n/open sum2slides --input meeting.txt --output slides.pptx\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🤖 与OpenClaw AI助手集成"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# AI助手自动调用\n助手：我可以帮您将文本转为幻灯片\n用户：好的，请处理这个会议纪要\n助手：正在使用sum2slides生成幻灯片...\n```"
    
    # 幻灯片5：OpenClaw配置文件集成
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    content = slide5.placeholders[1]
    
    title.text = "OpenClaw配置集成"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "⚙️ OpenClaw配置文件设置"
    p = text_frame.add_paragraph()
    p.text = "```yaml\n# ~/.openclaw/config.yaml\nskills:\n  sum2slides:\n    enabled: true\n    auto_trigger: true\n    default_template: business\n    output_dir: ~/Documents/Presentations\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🎯 环境变量配置"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 设置OpenClaw环境变量\nexport OPENCLAW_SUM2SLIDES_TEMPLATE=business\nexport OPENCLAW_SUM2SLIDES_THEME=light\nexport OPENCLAW_SUM2SLIDES_OUTPUT_DIR=~/presentations\n```"
    
    # 幻灯片6：OpenClaw工作流集成
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    content = slide6.placeholders[1]
    
    title.text = "OpenClaw工作流集成"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🔄 自动化工作流示例"
    p = text_frame.add_paragraph()
    p.text = "```yaml\n# OpenClaw工作流配置\nworkflows:\n  meeting_to_slides:\n    trigger: \"收到会议纪要文件\"\n    actions:\n      - \"使用sum2slides转换\"\n      - \"保存到共享文件夹\"\n      - \"发送通知到Slack\"\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🤖 AI助手自动化脚本"
    p = text_frame.add_paragraph()
    p.text = "```python\n# OpenClaw AI助手脚本\nfrom openclaw.assistant import Assistant\nfrom sum2slides import Sum2Slides\n\nassistant = Assistant()\nassistant.register_skill('sum2slides', Sum2Slides())\n\n# 自动处理文本消息\ndef handle_message(message):\n    if '会议纪要' in message:\n        slides = assistant.skills['sum2slides'].convert(message)\n        return f\"已生成幻灯片: {slides}\"\n```"
    
    # 幻灯片7：OpenClaw与其他Skill集成
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    content = slide7.placeholders[1]
    
    title.text = "与其他OpenClaw Skill集成"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🔗 技能组合使用"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 与文件管理Skill集成\nopenclaw skill install file-manager\n\n# 工作流：文件 → 转换 → 分享\n1. 使用file-manager获取文件\n2. 使用sum2slides转换\n3. 使用share-skill分享结果\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🤝 常用技能组合"
    p = text_frame.add_paragraph()
    p.text = "1. **sum2slides + calendar-skill**\n   会议安排 → 自动生成会议纪要幻灯片"
    p = text_frame.add_paragraph()
    p.text = "2. **sum2slides + email-skill**\n   邮件内容 → 转换为演示文稿"
    p = text_frame.add_paragraph()
    p.text = "3. **sum2slides + chatbot-skill**\n   聊天记录 → 整理为演示文稿"
    
    # 幻灯片8：OpenClaw使用场景
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    content = slide8.placeholders[1]
    
    title.text = "OpenClaw使用场景"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🎯 场景1：OpenClaw聊天机器人"
    p = text_frame.add_paragraph()
    p.text = "```\n用户：@sum2slides 请将今天的会议讨论转为幻灯片\n助手：好的，正在处理...\n      已生成幻灯片，请查收\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🎯 场景2：OpenClaw自动化工作流"
    p = text_frame.add_paragraph()
    p.text = "```yaml\n# 每天自动生成日报\nschedule:\n  daily_report:\n    time: \"17:00\"\n    actions:\n      - \"收集当日工作记录\"\n      - \"使用sum2slides生成幻灯片\"\n      - \"发送到团队频道\"\n```"
    
    p = text_frame.add_paragraph()
    p.text = "🎯 场景3：OpenClaw API集成"
    p = text_frame.add_paragraph()
    p.text = "```python\n# 在自定义应用中使用\nfrom openclaw_api import OpenClawClient\n\nclient = OpenClawClient()\nresult = client.call_skill('sum2slides', {\n    'text': '项目进度报告...',\n    'template': 'business'\n})\n```"
    
    # 幻灯片9：OpenClaw Skill管理
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide9.shapes.title
    content = slide9.placeholders[1]
    
    title.text = "OpenClaw Skill管理"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "🔄 技能管理命令"
    p = text_frame.add_paragraph()
    p.text = "```bash\n# 查看已安装技能\nopenclaw skill list\n\n# 查看sum2slides详情\nopenclaw skill info sum2slides\n\n# 更新sum2slides\nopenclaw skill update sum2slides\n\n# 禁用/启用技能\nopenclaw skill disable sum2slides\nopenclaw skill enable sum2slides\n\n# 卸载技能\nopenclaw skill uninstall sum2slides\n```"
    
    # 幻灯片10：OpenClaw最佳实践
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide10.shapes.title
    content = slide10.placeholders[1]
    
    title.text = "OpenClaw最佳实践"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "💡 实践建议"
    p = text_frame.add_paragraph()
    p.text = "1. **配置自动触发**：设置关键词自动调用sum2slides"
    p = text_frame.add_paragraph()
    p.text = "2. **预设模板**：为不同场景配置默认模板"
    p = text_frame.add_paragraph()
    p.text = "3. **输出目录**：设置统一的幻灯片输出目录"
    p = text_frame.add_paragraph()
    p.text = "4. **权限管理**：配置合适的文件访问权限"
    p = text_frame.add_paragraph()
    p.text = "5. **备份策略**：定期备份生成的演示文稿"
    
    # 幻灯片11：资源和支持
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide11.shapes.title
    content = slide11.placeholders[1]
    
    title.text = "OpenClaw Skill资源"
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "📚 OpenClaw相关文档"
    p = text_frame.add_paragraph()
    p.text = "  • OpenClaw官方文档: https://docs.openclaw.ai"
    p = text_frame.add_paragraph()
    p.text = "  • Skill开发指南: https://docs.openclaw.ai/skills"
    p = text_frame.add_paragraph()
    p.text = "  • API参考: https://docs.openclaw.ai/api"
    
    p = text_frame.add_paragraph()
    p.text = "🔗 社区和支持"
    p = text_frame.add_paragraph()
    p.text = "  • OpenClaw社区: https://discord.gg/clawd"
    p = text_frame.add_paragraph()
    p.text = "  • GitHub仓库: 本地化技能版本"
    p = text_frame.add_paragraph()
    p.text = "  • 问题反馈: 本地化技能版本/issues"
    
    # 幻灯片12：总结
    slide12 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide12.shapes.title
    subtitle = slide12.placeholders[1]
    
    title.text = "Sum2Slides for OpenClaw"
    subtitle.text = "专为OpenClaw设计的文本转幻灯片Skill\n\n🚀 一键安装，轻松使用\n🤖 深度集成，智能工作流\n💡 丰富场景，高效办公\n\n版本: Sum2Slides v1.0.1\n日期: " + datetime.now().strftime("%Y-%m-%d")
    
    # 保存文件
    output_dir = "openclaw_demo"
    os.makedirs(output_dir, exist_ok=True)
    
    output_file = f"{output_dir}/sum2slides_openclaw_demo.pptx"
    prs.save(output_file)
    
    print(f"✅ OpenClaw演示幻灯片已生成: {output_file}")
    print(f"📊 幻灯片数量: {len(prs.slides)} 张")
    
    return output_file

def main():
    """主函数"""
    
    print("🎬 生成Sum2Slides v1.0.1 OpenClaw专用演示幻灯片")
    print("=" * 50)
    
    try:
        # 检查python-pptx是否安装
        import pptx
        print("✅ python-pptx 已安装")
        
        # 生成演示文稿
        output_file = create_openclaw_demo_presentation()
        
        print("\n✅ OpenClaw演示幻灯片生成完成！")
        print(f"\n📁 输出文件: {output_file}")
        print("\n📋 幻灯片内容:")
        print("  1. 标题页 - Sum2Slides v1.0.1 (OpenClaw版)")
        print("  2. OpenClaw Skill定位")
        print("  3. OpenClaw安装方式")
        print("  4. OpenClaw环境使用")
        print("  5. OpenClaw配置集成")
        print("  6. OpenClaw工作流集成")
        print("  7. 与其他Skill集成")
        print("  8. OpenClaw使用场景")
        print("  9. OpenClaw Skill管理")
        print("  10. OpenClaw最佳实践")
        print("  11. OpenClaw资源支持")
        print("  12. 总结")
        
        print("\n🎯 重点内容:")
        print("  • 一键安装命令: openclaw skill install sum2slides")
        print("  • OpenClaw聊天使用: @sum2slides 或 /open sum2slides")
        print("  • 配置文件集成: ~/.openclaw/config.yaml")
        print("  • 自动化工作流: 会议纪要→幻灯片自动化")
        print("  • 技能组合: 与其他OpenClaw Skill深度集成")
        
        # 复制到samba共享文件夹
        samba_path = "~/samba_shares/小七文档/skill-versions/sum2slides_openclaw_demo.pptx"
        import shutil
        shutil.copy(output_file, os.path.expanduser(samba_path))
        
        print(f"\n📤 已复制到Samba共享文件夹: {samba_path}")
        
    except ImportError as e:
        print(f"❌ 需要安装python-pptx: {e}")
        print("\n💡 安装方法:")
        print("  pip install python-pptx")

if __name__ == "__main__":
    main()