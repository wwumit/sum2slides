"""
PowerPoint导出器
"""

import os
from typing import Optional, Dict, Any
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..models.slide import Presentation, Slide, Layout


class PowerPointExporter:
    """PowerPoint输出器"""
    
    def __init__(self):
        self.default_template = None
        self.layout_map = {
            Layout.TITLE: 0,                    # 标题幻灯片
            Layout.TITLE_AND_CONTENT: 1,        # 标题和内容
            Layout.SECTION_HEADER: 2,           # 章节标题
            Layout.TWO_CONTENT: 3,              # 两栏内容
            Layout.COMPARISON: 4,               # 比较
            Layout.TITLE_ONLY: 5,               # 仅标题
            Layout.BLANK: 6,                    # 空白
            Layout.CONTENT_WITH_CAPTION: 7,     # 带标题的内容
            Layout.PICTURE_WITH_CAPTION: 8,     # 带标题的图片
        }
    
    def export(self, presentation: Presentation, 
              output_path: str, options: Optional[Dict[str, Any]] = None) -> None:
        """导出为PowerPoint文件"""
        
        if options is None:
            options = {}
        
        # 创建PPTX演示文稿
        pptx_presentation = self._create_pptx_presentation(presentation, options)
        
        # 添加幻灯片
        for slide in presentation.slides:
            self._add_slide_to_pptx(pptx_presentation, slide, options)
        
        # 保存文件
        self._save_pptx(pptx_presentation, output_path)
    
    def set_options(self, options: Dict[str, Any]) -> None:
        """设置导出选项"""
        
        # 这里可以添加选项验证和设置逻辑
        self.options = options
    
    def _create_pptx_presentation(self, presentation: Presentation, 
                                 options: Dict[str, Any]) -> PPTXPresentation:
        """创建PPTX演示文稿对象"""
        
        # 使用默认模板或指定模板
        template_path = options.get('template_path')
        if template_path and os.path.exists(template_path):
            pptx_presentation = PPTXPresentation(template_path)
        else:
            pptx_presentation = PPTXPresentation()
        
        # 设置演示文稿属性
        pptx_presentation.core_properties.title = presentation.title
        pptx_presentation.core_properties.author = presentation.author or "Sum2Slides"
        pptx_presentation.core_properties.subject = "自动生成的演示文稿"
        pptx_presentation.core_properties.keywords = "Sum2Slides, 自动生成, 演示文稿"
        
        return pptx_presentation
    
    def _add_slide_to_pptx(self, pptx_presentation: PPTXPresentation, 
                          slide: Slide, options: Dict[str, Any]) -> None:
        """添加幻灯片到PPTX演示文稿"""
        
        # 确定幻灯片布局
        layout_index = self.layout_map.get(slide.layout, 1)  # 默认使用标题和内容布局
        
        # 创建幻灯片
        slide_layout = pptx_presentation.slide_layouts[layout_index]
        pptx_slide = pptx_presentation.slides.add_slide(slide_layout)
        
        # 设置幻灯片标题
        title_shape = pptx_slide.shapes.title
        if title_shape:
            title_shape.text = slide.title
        
        # 设置幻灯片内容
        if slide.layout in [Layout.TITLE_AND_CONTENT, Layout.TWO_CONTENT, 
                           Layout.COMPARISON, Layout.CONTENT_WITH_CAPTION]:
            self._add_content_to_slide(pptx_slide, slide, options)
        
        # 添加幻灯片编号
        self._add_slide_number(pptx_slide, slide.slide_number)
        
        # 添加演讲者备注
        if slide.notes:
            self._add_slide_notes(pptx_slide, slide.notes)
    
    def _add_content_to_slide(self, pptx_slide, slide: Slide, options: Dict[str, Any]) -> None:
        """添加内容到幻灯片"""
        
        # 获取内容占位符
        content_placeholder = None
        for shape in pptx_slide.placeholders:
            if shape.placeholder_format.type == 7:  # 正文内容
                content_placeholder = shape
                break
        
        if content_placeholder and slide.content:
            # 清空默认内容
            content_placeholder.text = ""
            
            # 添加内容
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # 清除所有段落
            
            # 设置文本格式
            for i, content in enumerate(slide.content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = content
                p.level = 0  # 一级列表
                
                # 设置字体
                font = p.font
                font_size = options.get('font_size', 18)
                font.size = Pt(font_size)
                font.name = options.get('font_name', 'Arial')
                
                # 设置颜色
                theme = options.get('theme', 'light')
                if theme == 'dark':
                    font.color.rgb = RGBColor(255, 255, 255)  # 白色
                else:
                    font.color.rgb = RGBColor(0, 0, 0)  # 黑色
    
    def _add_slide_number(self, pptx_slide, slide_number: int) -> None:
        """添加幻灯片编号"""
        
        # 在幻灯片底部添加编号
        # 注意：这需要幻灯片母版支持
        try:
            # 尝试使用页脚占位符
            for shape in pptx_slide.shapes:
                if shape.has_text_frame:
                    if "页脚" in shape.text or "Footer" in shape.text:
                        shape.text = f"{slide_number}"
                        break
        except:
            # 如果无法添加页脚，跳过
            pass
    
    def _add_slide_notes(self, pptx_slide, notes: str) -> None:
        """添加演讲者备注"""
        
        try:
            notes_slide = pptx_slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            
            # 清空现有备注
            notes_frame.clear()
            
            # 添加新备注
            p = notes_frame.paragraphs[0]
            p.text = notes
        except:
            # 如果无法添加备注，跳过
            pass
    
    def _save_pptx(self, pptx_presentation: PPTXPresentation, output_path: str) -> None:
        """保存PPTX文件"""
        
        # 确保目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        
        # 保存文件
        pptx_presentation.save(output_path)
    
    def _apply_theme(self, pptx_presentation: PPTXPresentation, theme: str) -> None:
        """应用主题"""
        
        # 这里可以添加主题应用逻辑
        # 注意：python-pptx的主题支持有限，可能需要手动设置颜色
        
        if theme == 'dark':
            # 设置深色主题颜色
            # 这里简化处理，实际需要设置更多属性
            pass
        elif theme == 'light':
            # 设置浅色主题颜色
            pass
        # 其他主题...